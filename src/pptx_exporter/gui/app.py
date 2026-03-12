"""App — main application window assembling all GUI components."""

import logging
import os
import subprocess
import sys
import threading
from pathlib import Path
from tkinter import filedialog, messagebox
from typing import Optional

import customtkinter as ctk

from ..exporter import Exporter
from ..utils import configure_logging, parse_slide_range
from .settings import load_settings, save_settings
from .tokens import COLORS, FONTS, RADIUS, SP, init_fonts
from .widgets import (
    ActionArea,
    ErrorBanner,
    FilePanel,
    SettingsCard,
    StatusPill,
)

logger = logging.getLogger(__name__)

ctk.set_default_color_theme("blue")


def _app_version() -> str:
    try:
        from importlib.metadata import version
        return version("pptx-exporter")
    except Exception:
        return ""


class App(ctk.CTk):
    """Main application window."""

    def __init__(self) -> None:
        super().__init__()
        init_fonts()
        configure_logging()

        self.title("pptx-exporter")
        self.resizable(True, True)
        self.minsize(560, 0)

        self._exporter = Exporter()
        self._pptx_paths: list[str] = []
        self._slide_counts: dict[str, int] = {}
        self._powerpoint_available = self._exporter.backend != "not_found"
        self._cancel_event: Optional[threading.Event] = None
        self._dnd_enabled = False

        settings = load_settings()
        self._ppi: int = settings.get("ppi", 300)
        self._theme = settings.get("theme", "light")
        ctk.set_appearance_mode(self._theme)
        saved_out = settings.get("output_dir")
        self._output_dir: Optional[str] = (
            saved_out if saved_out and os.path.isdir(saved_out) else None
        )

        self._build_ui()
        self._init_dnd()
        self._update_export_state()

        # Keyboard shortcuts
        self.bind("<Command-o>", lambda _: self._browse_pptx())
        self.bind("<Control-o>", lambda _: self._browse_pptx())
        self.bind("<Escape>", lambda _: self._on_escape())

    # ------------------------------------------------------------------
    # UI construction
    # ------------------------------------------------------------------

    def _build_ui(self) -> None:
        self.configure(fg_color=COLORS["bg"])
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(2, weight=1)  # content row expands

        # -- Toolbar -----------------------------------------------------------
        toolbar = ctk.CTkFrame(self, fg_color=COLORS["surface"], corner_radius=0)
        toolbar.grid(row=0, column=0, sticky="ew")
        toolbar.grid_columnconfigure(1, weight=1)

        # Left accent stripe
        ctk.CTkFrame(
            toolbar, width=4, fg_color=COLORS["accent"], corner_radius=0,
        ).grid(row=0, column=0, sticky="ns")

        # Title area
        title_frame = ctk.CTkFrame(toolbar, fg_color="transparent")
        title_frame.grid(row=0, column=1, sticky="w", padx=SP["md"], pady=SP["md"])

        ctk.CTkLabel(
            title_frame,
            text="PPTX EXPORTER",
            font=FONTS["title"],
            text_color=COLORS["text_primary"],
            anchor="w",
        ).grid(row=0, column=0, sticky="w")

        # Version badge
        ver = _app_version()
        if ver:
            ctk.CTkLabel(
                title_frame,
                text=f"v{ver}",
                font=FONTS["caption"],
                text_color=COLORS["text_tertiary"],
                anchor="w",
            ).grid(row=0, column=1, sticky="w", padx=(SP["sm"], 0))

        # Status pill
        self._status_pill = StatusPill(toolbar)
        self._status_pill.grid(
            row=0, column=2, sticky="e", padx=SP["md"], pady=SP["md"],
        )
        if self._powerpoint_available:
            self._status_pill.set_ready()
        else:
            self._status_pill.set_error()

        # Theme toggle
        theme_icon = "\u2600" if self._theme == "light" else "\u263D"
        self._theme_btn = ctk.CTkButton(
            toolbar,
            text=theme_icon,
            command=self._toggle_theme,
            width=32,
            height=32,
            font=(FONTS["body"][0], 16),
            fg_color="transparent",
            text_color=COLORS["text_secondary"],
            hover_color=COLORS["surface_hover"],
            corner_radius=RADIUS["sm"],
        )
        self._theme_btn.grid(
            row=0, column=3, sticky="e", padx=(0, SP["md"]), pady=SP["md"],
        )

        # Toolbar bottom border
        ctk.CTkFrame(
            self, height=1, fg_color=COLORS["border"], corner_radius=0,
        ).grid(row=1, column=0, sticky="ew")

        # -- Content area ------------------------------------------------------
        content = ctk.CTkFrame(self, fg_color="transparent")
        content.grid(
            row=2, column=0, sticky="nsew", padx=SP["lg"], pady=SP["lg"],
        )
        content.grid_columnconfigure(0, weight=1)
        content.grid_rowconfigure(0, weight=1)  # file panel grows

        # File panel (drop zone / file list) — expandable
        self._file_panel = FilePanel(
            content,
            on_browse=self._browse_pptx,
            on_clear_all=self._clear_pptx,
            on_remove_file=self._remove_file,
            dnd_enabled=self._dnd_enabled,
        )
        self._file_panel.grid(
            row=0, column=0, sticky="nsew", pady=(0, SP["md"]),
        )

        # Settings card — pinned to bottom
        self._settings_card = SettingsCard(
            content,
            initial_ppi=self._ppi,
            initial_output=self._output_dir,
            on_ppi_change=self._on_ppi_change,
            on_browse_output=self._browse_output,
            on_slide_toggle=self._on_slide_toggle,
        )
        self._settings_card.grid(row=1, column=0, sticky="sew", pady=(0, SP["md"]))

        # Action area — pinned to bottom
        self._action_area = ActionArea(
            content,
            on_run=self._on_run,
            on_cancel=self._on_cancel,
            on_open_folder=self._open_output_folder,
        )
        self._action_area.grid(row=2, column=0, sticky="sew", pady=(0, SP["sm"]))

        # Error banner (hidden) — pinned to bottom
        self._error_banner = ErrorBanner(content, on_dismiss=self._dismiss_error)
        self._error_banner.grid(row=3, column=0, sticky="sew")
        self._error_banner.grid_remove()

        # Set minimum height
        self.update_idletasks()
        self.minsize(560, self.winfo_reqheight())

    # ------------------------------------------------------------------
    # Drag-and-drop
    # ------------------------------------------------------------------

    def _init_dnd(self) -> None:
        try:
            from ..tkdnd import init_dnd, register_drop_target
            from ..tkdnd import bind_drop, bind_drop_enter, bind_drop_leave
            if not init_dnd(self):
                return
            register_drop_target(self._file_panel)
            bind_drop(self._file_panel, self._on_drop)
            bind_drop_enter(self._file_panel, self._on_drop_enter)
            bind_drop_leave(self._file_panel, self._on_drop_leave)
            self._dnd_enabled = True
            logger.debug("Drag-and-drop enabled")
        except Exception as exc:
            logger.debug("Drag-and-drop unavailable: %s", exc)

    def _on_drop(self, paths: tuple[str, ...]) -> None:
        pptx = [p for p in paths if p.lower().endswith(".pptx")]
        skipped = len(paths) - len(pptx)
        if skipped:
            self._show_inline_error(
                f"Only .pptx files are supported. "
                f"{skipped} file{'s' if skipped != 1 else ''} skipped."
            )
        if pptx:
            self._add_files(pptx)
        self._file_panel.set_drop_highlight(False)

    def _on_drop_enter(self) -> None:
        self._file_panel.set_drop_highlight(True)

    def _on_drop_leave(self) -> None:
        self._file_panel.set_drop_highlight(False)

    # ------------------------------------------------------------------
    # File management
    # ------------------------------------------------------------------

    def _browse_pptx(self) -> None:
        paths = filedialog.askopenfilenames(
            title="Select PowerPoint file(s)",
            filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")],
        )
        if paths:
            self._add_files(list(paths))

    def _add_files(self, new_paths: list[str]) -> None:
        dupes = [p for p in new_paths if p in self._pptx_paths]
        if dupes:
            names = ", ".join(Path(p).name for p in dupes)
            self._show_inline_error(f"Already added: {names}")
        unique = [p for p in new_paths if p not in self._pptx_paths]
        if not unique:
            return

        for path in unique:
            try:
                from pptx import Presentation
                self._slide_counts[path] = len(Presentation(path).slides)
            except Exception:
                self._slide_counts[path] = 0

        self._pptx_paths.extend(unique)
        logger.debug(
            "Added %d file(s), total now %d", len(unique), len(self._pptx_paths),
        )

        if self._output_dir is None:
            if len(self._pptx_paths) == 1:
                p = Path(self._pptx_paths[0])
                default_out = str(p.parent / f"{p.stem}_pngs")
            else:
                default_out = str(Path(self._pptx_paths[0]).parent)
            self._output_dir = default_out
            self._settings_card.set_output_dir(default_out)

        self._refresh_file_panel()
        self._update_slides_section()
        self._update_export_state()

    def _remove_file(self, path: str) -> None:
        if path in self._pptx_paths:
            self._pptx_paths.remove(path)
            self._slide_counts.pop(path, None)
        if not self._pptx_paths:
            self._clear_pptx()
            return
        self._refresh_file_panel()
        self._update_slides_section()
        self._update_export_state()

    def _clear_pptx(self) -> None:
        self._pptx_paths = []
        self._slide_counts.clear()
        self._file_panel.show_empty()
        self._settings_card.hide_slides()
        self._update_export_state()

    def _refresh_file_panel(self) -> None:
        file_infos = [
            (p, self._slide_counts.get(p, 0)) for p in self._pptx_paths
        ]
        self._file_panel.set_files(file_infos)

    def _update_slides_section(self) -> None:
        if (len(self._pptx_paths) == 1
                and self._slide_counts.get(self._pptx_paths[0], 0) > 1):
            self._settings_card.show_slides()
        else:
            self._settings_card.hide_slides()

    def _on_slide_toggle(self) -> None:
        self._settings_card.clear_slide_error()

    # ------------------------------------------------------------------
    # Settings
    # ------------------------------------------------------------------

    def _browse_output(self) -> None:
        path = filedialog.askdirectory(title="Select output folder")
        if path:
            self._output_dir = path
            self._settings_card.set_output_dir(path)
            logger.debug("Selected output dir: %s", path)
            self._save_settings()
            self._update_export_state()

    def _on_ppi_change(self, value: int) -> None:
        self._ppi = value
        logger.debug("PPI set to %d", self._ppi)
        self._save_settings()

    # ------------------------------------------------------------------
    # Export
    # ------------------------------------------------------------------

    def _on_run(self) -> None:
        if not self._pptx_paths or not self._output_dir:
            return

        # Parse slide selection (single-file only)
        self._slide_indices = None
        sc = self._slide_counts.get(self._pptx_paths[0], 0)
        if (len(self._pptx_paths) == 1
                and sc > 1
                and not self._settings_card.all_slides):
            spec = self._settings_card.get_slide_range_text()
            if not spec:
                self._settings_card.show_slide_error(
                    "Enter a slide range or check \"All\"",
                )
                return
            try:
                self._slide_indices = parse_slide_range(spec, sc)
                self._settings_card.clear_slide_error()
            except ValueError as exc:
                self._settings_card.show_slide_error(str(exc))
                return

        # Warn on overwrite
        out = Path(self._output_dir)
        glob_pattern = (
            "**/slide_*.png" if len(self._pptx_paths) > 1 else "slide_*.png"
        )
        if out.is_dir() and list(out.glob(glob_pattern)):
            ok = messagebox.askyesno(
                "Overwrite existing files?",
                f"The output folder already contains exported slides.\n\n"
                f"{self._output_dir}\n\n"
                "Existing slide PNGs will be overwritten. Continue?",
                parent=self,
            )
            if not ok:
                return

        self._error_banner.grid_remove()
        self._cancel_event = threading.Event()
        self._action_area.set_busy(True)
        thread = threading.Thread(target=self._run_export, daemon=True)
        thread.start()

    def _on_cancel(self) -> None:
        if self._cancel_event:
            self._cancel_event.set()
        self._action_area.set_cancelling()

    def _on_escape(self) -> None:
        if self._cancel_event and not self._cancel_event.is_set():
            self._on_cancel()

    def _open_output_folder(self) -> None:
        if not self._output_dir:
            return
        if sys.platform == "darwin":
            subprocess.run(["open", self._output_dir], check=False)
        elif sys.platform == "win32":
            os.startfile(self._output_dir)  # type: ignore[attr-defined]

    # ------------------------------------------------------------------
    # Export thread
    # ------------------------------------------------------------------

    def _run_export(self) -> None:
        try:
            paths = list(self._pptx_paths)
            if len(paths) == 1:
                self._exporter.export(
                    paths[0],
                    self._output_dir,
                    progress_callback=self._on_progress,
                    cancel_event=self._cancel_event,
                    ppi=self._ppi,
                    slide_indices=self._slide_indices,
                )
            else:
                self._run_batch_export(paths)
            self.after(0, self._on_export_done)
        except InterruptedError:
            self.after(0, self._on_export_cancelled)
        except Exception as exc:
            logger.exception("Export failed")
            msg = str(exc)
            if self._output_dir:
                exported = list(Path(self._output_dir).glob(
                    "**/slide_*.png" if len(self._pptx_paths) > 1
                    else "slide_*.png"
                ))
                if exported:
                    msg += (
                        f"\n\n{len(exported)} slide(s) were exported "
                        "before the error."
                    )
            self.after(0, self._on_export_error, msg)

    def _run_batch_export(self, paths: list[str]) -> None:
        from pptx import Presentation as PptxPresentation

        slide_counts = []
        for p in paths:
            try:
                slide_counts.append(len(PptxPresentation(p).slides))
            except Exception:
                slide_counts.append(1)
        grand_total = sum(slide_counts)

        for file_idx, pptx_path in enumerate(paths):
            if self._cancel_event and self._cancel_event.is_set():
                raise InterruptedError("Export cancelled by user.")

            stem = Path(pptx_path).stem
            file_out = str(Path(self._output_dir) / f"{stem}_pngs")
            offset = sum(slide_counts[:file_idx])

            def _make_cb(off, gt, fidx, fname, fcount):
                def cb(current, total):
                    abs_current = off + current
                    frac = abs_current / gt if gt > 0 else 0
                    if current < total:
                        msg = (
                            f"[{fidx}/{fcount}] {fname}: "
                            f"slide {current + 1}/{total}"
                        )
                    else:
                        msg = f"[{fidx}/{fcount}] {fname}: done"
                    self.after(0, self._action_area.update_progress, frac, msg)
                return cb

            self._exporter.export(
                pptx_path,
                file_out,
                progress_callback=_make_cb(
                    offset, grand_total, file_idx + 1, stem, len(paths),
                ),
                cancel_event=self._cancel_event,
                ppi=self._ppi,
            )

    def _on_progress(self, current: int, total: int) -> None:
        if total == 0:
            return
        fraction = current / total
        msg = (
            f"Processing slide {current + 1} of {total}\u2026"
            if current < total
            else "Finalising\u2026"
        )
        self.after(0, self._action_area.update_progress, fraction, msg)

    def _on_export_done(self) -> None:
        n = len(self._pptx_paths)
        if n > 1:
            msg = f"Done \u2014 {n} presentations exported."
        else:
            msg = "Done \u2014 all slides exported."
        self._action_area.set_busy(False)
        self._action_area.show_done(msg)
        self._update_export_state()

    def _on_export_cancelled(self) -> None:
        self._action_area.set_busy(False)
        self._action_area.show_cancelled()
        self._update_export_state()

    def _on_export_error(self, message: str) -> None:
        self._action_area.set_busy(False)
        self._update_export_state()
        self._error_banner.show(message)
        self._error_banner.grid()

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _toggle_theme(self) -> None:
        self._theme = "dark" if self._theme == "light" else "light"
        ctk.set_appearance_mode(self._theme)
        icon = "\u2600" if self._theme == "light" else "\u263D"
        self._theme_btn.configure(text=icon)
        self._save_settings()

    def _save_settings(self) -> None:
        save_settings({
            "ppi": self._ppi,
            "output_dir": self._output_dir,
            "theme": self._theme,
        })

    def _update_export_state(self) -> None:
        ready = bool(
            self._powerpoint_available and self._pptx_paths and self._output_dir
        )
        if not self._powerpoint_available:
            hint = "PowerPoint not found"
        elif not self._pptx_paths:
            hint = "Select a file to begin"
        elif not self._output_dir:
            hint = "Choose an output folder"
        else:
            hint = ""
        self._action_area.set_ready(ready, hint)

    def _dismiss_error(self) -> None:
        self._error_banner.grid_remove()

    def _show_inline_error(self, msg: str) -> None:
        """Show a transient error via the error banner."""
        self._error_banner.show(msg)
        self._error_banner.grid()

    def _set_ui_busy(self, busy: bool) -> None:
        self._action_area.set_busy(busy)
        if not busy:
            self._update_export_state()
        self.update_idletasks()
