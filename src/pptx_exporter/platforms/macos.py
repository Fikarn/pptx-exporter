"""macOS backend: drive Microsoft PowerPoint via AppleScript / osascript.

Workflow:
1. Open the original .pptx in PowerPoint.
2. For each slide:
   a. Add an invisible full-slide rectangle via AppleScript.
   b. Select all shapes (Esc, Esc, Cmd+A).
   c. Copy the selection (Cmd+C) — this captures all objects including
      placeholders as a transparent PNG on the clipboard.
   d. Save the clipboard PNG to disk via NSPasteboard.
3. Close without saving — the original file is never modified.

The clipboard approach is necessary because AppleScript's ``save as picture``
only works on a single shape, and PowerPoint refuses to group placeholder
shapes. Copying a multi-selection preserves transparency and includes all
objects.
"""

import logging
import subprocess
import threading
from pathlib import Path
from typing import Callable, Optional, Sequence

from ..utils import slide_output_name

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# AppleScript templates
# ---------------------------------------------------------------------------

_SCRIPT_OPEN = """\
tell application "Microsoft PowerPoint"
    activate
    open POSIX file "{pptx_path}"
    -- Poll until the presentation is fully loaded (up to 15 s).
    repeat 30 times
        try
            if (count of slides of active presentation) > 0 then return "ready"
        end try
        delay 0.5
    end repeat
    return "timeout"
end tell
"""

_SCRIPT_GO_TO_SLIDE = """\
tell application "Microsoft PowerPoint"
    activate
    set theView to view of active window
    set slide of theView to slide {slide_num} of active presentation
    -- Poll until navigation completes (up to 5 s).
    repeat 20 times
        try
            if (slide number of slide of theView) = {slide_num} then return "ready"
        end try
        delay 0.25
    end repeat
    return "timeout"
end tell
"""

# Add an invisible full-slide rectangle to the current slide.
_SCRIPT_ADD_BOUNDING_RECT = """\
tell application "Microsoft PowerPoint"
    set oSlide to slide {slide_num} of active presentation
    set newShape to make new shape at oSlide with properties {{auto shape type:autoshape rectangle, left position:0, top:0, width:{slide_w}, height:{slide_h}}}  # noqa: E501
    set visible of fill format of newShape to false
    set line weight of line format of newShape to 0
    set transparency of line format of newShape to 1.0
    return "ok"
end tell
"""

# Select all shapes on the current slide and copy to clipboard.
_SCRIPT_SELECT_ALL_AND_COPY = """\
tell application "Microsoft PowerPoint" to activate
delay 0.3
tell application "System Events"
    tell process "Microsoft PowerPoint"
        key code 53
        delay 0.2
        key code 53
        delay 0.2
        keystroke "a" using {command down}
        delay 0.3
        keystroke "c" using {command down}
        delay 0.5
    end tell
end tell
"""

# Save clipboard image data to a high-resolution PNG via NSPasteboard (Cocoa).
# Tries PDF (vector) first for best quality, then TIFF, then PNG.
# Renders into a bitmap at the target pixel dimensions for 300 PPI output.
_SCRIPT_CLIPBOARD_TO_PNG = """\
use framework "AppKit"
use scripting additions

set pb to current application's NSPasteboard's generalPasteboard()
set targetW to {target_w} as integer
set targetH to {target_h} as integer

-- Collect clipboard data: prefer PDF (vector), then TIFF, then PNG
set srcData to missing value
set pdfType to current application's NSPasteboardTypePDF
set srcData to pb's dataForType:pdfType
if srcData is missing value then
    set tiffType to current application's NSPasteboardTypeTIFF
    set srcData to pb's dataForType:tiffType
end if
if srcData is missing value then
    set pngType to current application's NSPasteboardTypePNG
    set srcData to pb's dataForType:pngType
end if
if srcData is missing value then return "no_image"

-- Build NSImage from clipboard data
set img to current application's NSImage's alloc()'s initWithData:srcData
if img is missing value then return "no_image"

-- Create a bitmap at the target resolution
set bitmapRep to (current application's NSBitmapImageRep's alloc()'s initWithBitmapDataPlanes:(missing value) pixelsWide:targetW pixelsHigh:targetH bitsPerSample:8 samplesPerPixel:4 hasAlpha:true isPlanar:false colorSpaceName:(current application's NSCalibratedRGBColorSpace) bytesPerRow:0 bitsPerPixel:0)  # noqa: E501
bitmapRep's setSize:(current application's NSMakeSize(targetW, targetH))

-- Render the image into the bitmap at target size
current application's NSGraphicsContext's saveGraphicsState()
set ctx to (current application's NSGraphicsContext's graphicsContextWithBitmapImageRep:bitmapRep)
current application's NSGraphicsContext's setCurrentContext:ctx
ctx's setImageInterpolation:(current application's NSImageInterpolationHigh)
img's drawInRect:(current application's NSMakeRect(0, 0, targetW, targetH))
current application's NSGraphicsContext's restoreGraphicsState()

-- Save as PNG
set pngData to (bitmapRep's representationUsingType:(current application's NSBitmapImageFileTypePNG) |properties|:(missing value))  # noqa: E501
if pngData is missing value then return "convert_failed"

set writeResult to (pngData's writeToFile:"{out_path}" atomically:true) as boolean
if writeResult then return "ok"
return "write_failed"
"""

_SCRIPT_CLOSE_WITHOUT_SAVING = """\
tell application "Microsoft PowerPoint"
    close active presentation saving no
end tell
"""


def _check_accessibility() -> None:
    """Raise RuntimeError if the app lacks Accessibility (System Events) access.

    The export relies on System Events keystrokes (Esc, Cmd+A, Cmd+C).
    Without Accessibility permission these silently fail, producing empty
    clipboard reads for every slide.
    """
    script = (
        'use framework "ApplicationServices"\n'
        "return (current application's AXIsProcessTrusted()) as boolean"
    )
    try:
        result = subprocess.run(
            ["osascript", "-e", script],
            capture_output=True, text=True, check=False,
        )
        if result.stdout.strip() == "false":
            raise RuntimeError(
                "Accessibility access is required.\n"
                "Go to System Settings → Privacy & Security → Accessibility "
                "and enable access for this app, then try again."
            )
    except RuntimeError:
        raise
    except Exception as exc:
        # If the check itself fails, log and continue — don't block the
        # export on an optional pre-flight check.
        logger.debug("Accessibility pre-check failed: %s", exc)


def _escape_applescript_string(s: str) -> str:
    """Escape a string for safe embedding in AppleScript double-quoted literals."""
    return s.replace("\\", "\\\\").replace('"', '\\"')


def _run_applescript(script: str, timeout: int = 30) -> str:
    """Execute *script* via osascript and return stdout stripped.

    Args:
        timeout: Maximum seconds to wait before killing the process.
    """
    try:
        result = subprocess.run(
            ["osascript", "-e", script],
            capture_output=True,
            text=True,
            check=False,
            timeout=timeout,
        )
    except subprocess.TimeoutExpired:
        raise RuntimeError(
            f"AppleScript timed out after {timeout}s. "
            "PowerPoint may be unresponsive — try restarting it."
        )
    if result.returncode != 0:
        raise RuntimeError(
            f"AppleScript error (exit {result.returncode}):\n"
            f"STDOUT: {result.stdout.strip()}\n"
            f"STDERR: {result.stderr.strip()}"
        )
    return result.stdout.strip()


def export_slides(
    pptx_path: Path,
    output_dir: Path,
    progress_callback: Optional[Callable[[int, int], None]] = None,
    cancel_event: Optional[threading.Event] = None,
    ppi: int = 300,
    slide_indices: Optional[Sequence[int]] = None,
) -> None:
    """Export every slide of *pptx_path* as a transparent PNG into *output_dir*.

    Opens the original file in PowerPoint. For each slide: adds an invisible
    bounding rectangle, selects all, copies to clipboard, and saves the
    clipboard PNG via NSPasteboard. Closes without saving.

    Args:
        cancel_event: Optional threading.Event; if set between slides the
            export is aborted cleanly and the presentation is closed.
    """
    _check_accessibility()

    # Read slide metadata via python-pptx (no modification).
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    slide_w = round(prs.slide_width / 12700, 1)   # EMU → points
    slide_h = round(prs.slide_height / 12700, 1)   # EMU → points
    # Target pixel dimensions at the requested PPI (points are 1/72 inch).
    target_ppi = ppi
    target_w = int(round(slide_w / 72 * target_ppi))
    target_h = int(round(slide_h / 72 * target_ppi))
    logger.info("Slide count: %d, dimensions: %.0f x %.0f pts, target: %d x %d px (%d PPI)",
                total, slide_w, slide_h, target_w, target_h, target_ppi)

    # Open the original file in PowerPoint.
    pptx_posix = str(pptx_path)
    logger.info("Opening presentation in PowerPoint: %s", pptx_posix)
    open_result = _run_applescript(_SCRIPT_OPEN.format(
        pptx_path=_escape_applescript_string(pptx_posix)
    ), timeout=60)
    if open_result == "timeout":
        raise RuntimeError(
            "Timed out waiting for PowerPoint to open the presentation. "
            "Try closing other presentations and retrying."
        )

    try:
        indices = list(slide_indices) if slide_indices is not None else list(range(total))
        selected_total = len(indices)
        for progress_idx, idx in enumerate(indices):
            if cancel_event and cancel_event.is_set():
                logger.info("Export cancelled before slide %d", idx + 1)
                raise InterruptedError("Export cancelled by user.")

            slide_num = idx + 1
            if progress_callback:
                progress_callback(progress_idx, selected_total)

            logger.debug("Processing slide %d/%d", slide_num, total)

            out_name = slide_output_name(idx, total)
            out_path = output_dir / out_name

            # Step 1: Navigate to the slide.
            nav_result = _run_applescript(
                _SCRIPT_GO_TO_SLIDE.format(slide_num=slide_num)
            )
            if nav_result == "timeout":
                logger.warning("Slide %d: navigation timed out, proceeding anyway", slide_num)

            # Step 2: Add invisible full-slide bounding rectangle.
            _run_applescript(
                _SCRIPT_ADD_BOUNDING_RECT.format(
                    slide_num=slide_num,
                    slide_w=slide_w,
                    slide_h=slide_h,
                )
            )

            # Step 3+4: Select all, copy, and save clipboard as PNG.
            # Retry up to 3 times — the clipboard may not be ready on the
            # first attempt due to timing between the copy and the read.
            max_attempts = 3
            save_result = "no_image"
            for attempt in range(1, max_attempts + 1):
                _run_applescript(_SCRIPT_SELECT_ALL_AND_COPY)
                save_result = _run_applescript(
                    _SCRIPT_CLIPBOARD_TO_PNG.format(
                        out_path=_escape_applescript_string(str(out_path)),
                        target_w=target_w,
                        target_h=target_h,
                    )
                )
                if save_result == "ok":
                    break
                if attempt < max_attempts:
                    logger.debug(
                        "Slide %d: clipboard attempt %d/%d returned '%s', retrying",
                        slide_num, attempt, max_attempts, save_result,
                    )

            if save_result == "ok":
                logger.info("Exported slide %d → %s", slide_num, out_path)
            elif save_result == "no_image":
                logger.warning("Slide %d: no image data on clipboard after %d attempts",
                               slide_num, max_attempts)
            else:
                logger.warning("Slide %d: clipboard save returned '%s'", slide_num, save_result)

    finally:
        try:
            _run_applescript(_SCRIPT_CLOSE_WITHOUT_SAVING)
            logger.info("Closed presentation without saving")
        except Exception as exc:
            logger.warning("Could not close presentation: %s", exc)

    if progress_callback:
        progress_callback(selected_total, selected_total)
