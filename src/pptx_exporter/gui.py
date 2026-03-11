"""All UI code for pptx-exporter — built with CustomTkinter."""

import json
import logging
import os
import subprocess
import sys
import threading
import tkinter as tk
from pathlib import Path
from tkinter import filedialog, messagebox
from typing import Optional

import customtkinter as ctk

from .exporter import Exporter
from .utils import configure_logging, parse_slide_range

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Appearance
# ---------------------------------------------------------------------------

ctk.set_appearance_mode("system")       # follows OS light / dark setting
ctk.set_default_color_theme("blue")

# Design tokens — (light, dark)
_COLOR_BG_SURFACE = ("#FFFFFF", "#2C2C2E")
_COLOR_BG_BASE = ("#F5F5F7", "#1C1C1E")
_COLOR_BORDER = ("#D1D1D6", "#3A3A3C")
_COLOR_TEXT_PRIMARY = ("#1D1D1F", "#F5F5F7")
_COLOR_TEXT_SECONDARY = ("#86868B", "#8E8E93")
_COLOR_ACCENT = ("#0071E3", "#0A84FF")
_COLOR_ACCENT_HOVER = ("#0077ED", "#409CFF")
_COLOR_CANCEL = ("#636366", "#636366")
_COLOR_CANCEL_HOVER = ("#48484A", "#48484A")
_COLOR_ERROR_BG = ("#FEE2E2", "#3B0A0A")
_COLOR_ERROR_TEXT = ("#991B1B", "#FCA5A5")
_COLOR_PILL_READY_BG = ("#D1FAE5", "#052E16")
_COLOR_PILL_READY_TEXT = ("#065F46", "#6EE7B7")
_COLOR_PILL_ERROR_BG = ("#FEE2E2", "#3B0A0A")
_COLOR_PILL_ERROR_TEXT = ("#991B1B", "#FCA5A5")

# System font — renders correctly on both macOS (SF Pro) and Windows (Segoe UI)
_FONT_BODY = ("system-ui", 13)
_FONT_BODY_BOLD = ("system-ui", 13, "bold")
_FONT_SMALL = ("system-ui", 11)
_FONT_TITLE = ("system-ui", 15, "bold")

_PPI_PRESETS = {
    "72 dpi": 72,
    "150 dpi": 150,
    "300 dpi": 300,
}
_PPI_SEGMENT_VALUES = ["72 dpi", "150 dpi", "300 dpi", "Custom"]
_PPI_MIN = 36
_PPI_MAX = 2400

# ---------------------------------------------------------------------------
# Settings persistence (~/.pptx-exporter-settings.json)
# ---------------------------------------------------------------------------

_SETTINGS_PATH = Path.home() / ".pptx-exporter-settings.json"


def _load_settings() -> dict:
    try:
        with open(_SETTINGS_PATH) as fh:
            return json.load(fh)
    except Exception:
        return {}


def _save_settings(data: dict) -> None:
    try:
        with open(_SETTINGS_PATH, "w") as fh:
            json.dump(data, fh, indent=2)
    except Exception as exc:
        logger.debug("Could not save settings: %s", exc)


class App(ctk.CTk):
    """Main application window."""

    def __init__(self) -> None:
        super().__init__()
        configure_logging()

        self.title("pptx-exporter")
        self.resizable(True, False)
        self.minsize(520, 0)

        self._exporter = Exporter()
        self._pptx_paths: list[str] = []
        self._slide_counts: dict[str, int] = {}
        self._powerpoint_available = self._exporter.backend != "not_found"
        self._cancel_event: Optional[threading.Event] = None
        self._dnd_enabled = False

        settings = _load_settings()
        self._ppi: int = settings.get("ppi", 300)
        saved_out = settings.get("output_dir")
        self._output_dir: Optional[str] = (
            saved_out if saved_out and os.path.isdir(saved_out) else None
        )

        self._build_ui()
        self._init_dnd()
        self._update_run_button_state()

    # ------------------------------------------------------------------
    # UI construction
    # ------------------------------------------------------------------

    def _build_ui(self) -> None:
        PAD = 16

        self.grid_columnconfigure(0, weight=1)

        # ── Header ────────────────────────────────────────────────────
        header = ctk.CTkFrame(self, fg_color=_COLOR_BG_SURFACE, corner_radius=0)
        header.grid(row=0, column=0, sticky="ew", padx=0, pady=0)
        header.grid_columnconfigure(0, weight=1)

        ctk.CTkLabel(
            header,
            text="pptx-exporter",
            font=_FONT_TITLE,
            text_color=_COLOR_TEXT_PRIMARY,
            anchor="w",
        ).grid(row=0, column=0, sticky="w", padx=PAD, pady=(PAD, PAD))

        self._backend_pill = ctk.CTkLabel(
            header,
            text="",
            font=_FONT_SMALL,
            corner_radius=10,
            padx=10,
            pady=3,
        )
        self._backend_pill.grid(row=0, column=1, sticky="e", padx=PAD, pady=PAD)
        self._refresh_backend_pill()

        # Separator
        ctk.CTkFrame(self, height=1, fg_color=_COLOR_BORDER, corner_radius=0).grid(
            row=1, column=0, sticky="ew"
        )

        # ── Main content area ─────────────────────────────────────────
        content = ctk.CTkFrame(self, fg_color=_COLOR_BG_BASE, corner_radius=0)
        content.grid(row=2, column=0, sticky="ew")
        content.grid_columnconfigure(0, weight=1)

        # Input files section
        ctk.CTkLabel(
            content,
            text="INPUT FILES",
            font=_FONT_SMALL,
            text_color=_COLOR_TEXT_SECONDARY,
            anchor="w",
        ).grid(row=0, column=0, sticky="w", padx=PAD, pady=(PAD, 4))

        self._file_panel = _FilePanel(
            content,
            on_browse=self._browse_pptx,
            on_clear_all=self._clear_pptx,
            on_remove_file=self._remove_file,
            dnd_enabled=self._dnd_enabled,
        )
        self._file_panel.grid(row=1, column=0, sticky="ew", padx=PAD, pady=(0, PAD))

        # Separator
        ctk.CTkFrame(content, height=1, fg_color=_COLOR_BORDER, corner_radius=0).grid(
            row=2, column=0, sticky="ew", padx=PAD
        )

        # Slide selection section (hidden until a file is selected)
        self._slides_frame = ctk.CTkFrame(content, fg_color="transparent")
        self._slides_frame.grid(row=3, column=0, sticky="ew")
        self._slides_frame.grid_columnconfigure(0, weight=1)
        self._slides_frame.grid_remove()

        ctk.CTkLabel(
            self._slides_frame,
            text="SLIDES",
            font=_FONT_SMALL,
            text_color=_COLOR_TEXT_SECONDARY,
            anchor="w",
        ).grid(row=0, column=0, sticky="w", padx=PAD, pady=(PAD, 4))

        slide_row = ctk.CTkFrame(self._slides_frame, fg_color="transparent")
        slide_row.grid(row=1, column=0, sticky="ew", padx=PAD, pady=(0, 4))
        slide_row.grid_columnconfigure(1, weight=1)

        self._all_slides_var = tk.BooleanVar(value=True)
        self._all_slides_cb = ctk.CTkCheckBox(
            slide_row,
            text="All slides",
            variable=self._all_slides_var,
            command=self._on_slide_selection_toggle,
            font=_FONT_BODY,
        )
        self._all_slides_cb.grid(row=0, column=0, sticky="w")

        self._slide_range_entry = ctk.CTkEntry(
            slide_row, font=_FONT_BODY,
            placeholder_text="e.g. 1-5, 8, 10-12",
        )
        self._slide_range_entry.grid(row=0, column=1, sticky="ew", padx=(12, 0))
        self._slide_range_entry.grid_remove()

        # Separator after slides
        ctk.CTkFrame(self._slides_frame, height=1, fg_color=_COLOR_BORDER,
                     corner_radius=0).grid(
            row=2, column=0, sticky="ew", padx=PAD, pady=(PAD, 0)
        )

        # Separator (when slides section is hidden)
        self._sep_after_file = ctk.CTkFrame(
            content, height=1, fg_color=_COLOR_BORDER, corner_radius=0
        )

        # Output folder section
        ctk.CTkLabel(
            content,
            text="OUTPUT FOLDER",
            font=_FONT_SMALL,
            text_color=_COLOR_TEXT_SECONDARY,
            anchor="w",
        ).grid(row=4, column=0, sticky="w", padx=PAD, pady=(PAD, 4))

        out_row = ctk.CTkFrame(content, fg_color="transparent")
        out_row.grid(row=5, column=0, sticky="ew", padx=PAD, pady=(0, PAD))
        out_row.grid_columnconfigure(0, weight=1)

        self._out_var = tk.StringVar(value=self._output_dir or "(none selected)")
        ctk.CTkLabel(
            out_row,
            textvariable=self._out_var,
            font=_FONT_BODY,
            text_color=_COLOR_TEXT_SECONDARY,
            anchor="w",
        ).grid(row=0, column=0, sticky="ew", padx=(0, 8))

        ctk.CTkButton(
            out_row,
            text="Browse…",
            command=self._browse_output,
            width=90,
            font=_FONT_BODY,
            fg_color=_COLOR_BG_SURFACE,
            text_color=_COLOR_TEXT_PRIMARY,
            border_width=1,
            border_color=_COLOR_BORDER,
            hover_color=_COLOR_BG_BASE,
        ).grid(row=0, column=1)

        # Separator
        ctk.CTkFrame(content, height=1, fg_color=_COLOR_BORDER, corner_radius=0).grid(
            row=6, column=0, sticky="ew", padx=PAD
        )

        # Resolution section
        ctk.CTkLabel(
            content,
            text="RESOLUTION",
            font=_FONT_SMALL,
            text_color=_COLOR_TEXT_SECONDARY,
            anchor="w",
        ).grid(row=7, column=0, sticky="w", padx=PAD, pady=(PAD, 4))

        # Determine initial segment selection.
        initial_label = next(
            (lbl for lbl, val in _PPI_PRESETS.items() if val == self._ppi),
            "Custom",
        )
        self._ppi_seg = ctk.CTkSegmentedButton(
            content,
            values=_PPI_SEGMENT_VALUES,
            command=self._on_ppi_change,
            font=_FONT_BODY,
        )
        self._ppi_seg.set(initial_label)
        self._ppi_seg.grid(row=8, column=0, sticky="w", padx=PAD, pady=(0, 4))

        # Custom PPI entry (shown only when "Custom" is selected).
        self._custom_ppi_frame = ctk.CTkFrame(content, fg_color="transparent")
        self._custom_ppi_frame.grid(row=9, column=0, sticky="w", padx=PAD, pady=(0, PAD))

        self._custom_ppi_entry = ctk.CTkEntry(
            self._custom_ppi_frame, width=80, font=_FONT_BODY,
            placeholder_text="dpi",
        )
        self._custom_ppi_entry.grid(row=0, column=0)
        self._custom_ppi_entry.bind("<Return>", lambda _: self._apply_custom_ppi())
        self._custom_ppi_entry.bind("<FocusOut>", lambda _: self._apply_custom_ppi())

        self._custom_ppi_hint = ctk.CTkLabel(
            self._custom_ppi_frame,
            text=f"  {_PPI_MIN}–{_PPI_MAX}",
            font=_FONT_SMALL,
            text_color=_COLOR_TEXT_SECONDARY,
        )
        self._custom_ppi_hint.grid(row=0, column=1)

        if initial_label == "Custom":
            self._custom_ppi_entry.insert(0, str(self._ppi))
        else:
            self._custom_ppi_frame.grid_remove()

        # Separator
        ctk.CTkFrame(self, height=1, fg_color=_COLOR_BORDER, corner_radius=0).grid(
            row=3, column=0, sticky="ew"
        )

        # ── Footer (progress + action) ────────────────────────────────
        footer = ctk.CTkFrame(self, fg_color=_COLOR_BG_SURFACE, corner_radius=0)
        footer.grid(row=4, column=0, sticky="ew")
        footer.grid_columnconfigure(0, weight=1)

        # Progress area (hidden until export starts)
        self._progress_frame = ctk.CTkFrame(footer, fg_color="transparent")
        self._progress_frame.grid(
            row=0, column=0, columnspan=2, sticky="ew", padx=PAD, pady=(PAD, 0)
        )
        self._progress_frame.grid_columnconfigure(0, weight=1)
        self._progress_frame.grid_remove()

        self._progress_bar = ctk.CTkProgressBar(
            self._progress_frame,
            mode="determinate",
            height=6,
        )
        self._progress_bar.set(0)
        self._progress_bar.grid(row=0, column=0, sticky="ew", pady=(0, 6))

        self._status_var = tk.StringVar(value="")
        self._status_label = ctk.CTkLabel(
            self._progress_frame,
            textvariable=self._status_var,
            font=_FONT_SMALL,
            text_color=_COLOR_TEXT_SECONDARY,
            anchor="w",
        )
        self._status_label.grid(row=1, column=0, sticky="w")

        # Open folder button (right of status, hidden until export done)
        self._open_folder_btn = ctk.CTkButton(
            self._progress_frame,
            text="Open Folder ↗",
            command=self._open_output_folder,
            width=110,
            height=26,
            font=_FONT_SMALL,
            fg_color="transparent",
            text_color=_COLOR_ACCENT,
            hover_color=_COLOR_BG_BASE,
            border_width=1,
            border_color=_COLOR_ACCENT,
            corner_radius=4,
        )
        self._open_folder_btn.grid(row=1, column=1, padx=(8, 0))
        self._open_folder_btn.grid_remove()

        # Action buttons row
        btn_row = ctk.CTkFrame(footer, fg_color="transparent")
        btn_row.grid(row=1, column=0, sticky="ew", padx=PAD, pady=PAD)
        btn_row.grid_columnconfigure(0, weight=1)

        self._run_btn = ctk.CTkButton(
            btn_row,
            text="Export PNGs",
            command=self._on_run,
            font=_FONT_BODY_BOLD,
            height=40,
            corner_radius=8,
            fg_color=_COLOR_ACCENT,
            hover_color=_COLOR_ACCENT_HOVER,
        )
        self._run_btn.grid(row=0, column=0, sticky="ew")

        self._cancel_btn = ctk.CTkButton(
            btn_row,
            text="Cancel",
            command=self._on_cancel,
            font=_FONT_BODY_BOLD,
            height=40,
            corner_radius=8,
            fg_color=_COLOR_CANCEL,
            hover_color=_COLOR_CANCEL_HOVER,
        )
        self._cancel_btn.grid(row=0, column=0, sticky="ew")
        self._cancel_btn.grid_remove()

        # Error banner (hidden until an error occurs)
        self._error_banner = _ErrorBanner(footer, on_dismiss=self._dismiss_error)
        self._error_banner.grid(row=2, column=0, sticky="ew", padx=PAD, pady=(0, PAD))
        self._error_banner.grid_remove()

        # Final geometry update
        self.update_idletasks()
        self.minsize(520, self.winfo_reqheight())

    # ------------------------------------------------------------------
    # Backend pill
    # ------------------------------------------------------------------

    def _refresh_backend_pill(self) -> None:
        if self._powerpoint_available:
            self._backend_pill.configure(
                text="● PowerPoint ready",
                fg_color=_COLOR_PILL_READY_BG,
                text_color=_COLOR_PILL_READY_TEXT,
            )
        else:
            self._backend_pill.configure(
                text="● PowerPoint not found",
                fg_color=_COLOR_PILL_ERROR_BG,
                text_color=_COLOR_PILL_ERROR_TEXT,
            )

    # ------------------------------------------------------------------
    # Drag-and-drop initialisation
    # ------------------------------------------------------------------

    def _init_dnd(self) -> None:
        """Try to load tkdnd and register the file panel as a drop target."""
        try:
            from .tkdnd import init_dnd, register_drop_target
            from .tkdnd import bind_drop, bind_drop_enter, bind_drop_leave
            if not init_dnd(self):
                return
            register_drop_target(self._file_panel)
            bind_drop(self._file_panel, self._on_drop)
            bind_drop_enter(self._file_panel, self._on_drop_enter)
            bind_drop_leave(self._file_panel, self._on_drop_leave)
            self._dnd_enabled = True
            logger.debug("Drag-and-drop enabled")
        except Exception as exc:
            logger.debug("Drag-and-drop unavailable: %s", exc)

    def _on_drop(self, paths: tuple[str, ...]) -> None:
        """Handle files dropped onto the file panel."""
        pptx = [p for p in paths if p.lower().endswith(".pptx")]
        skipped = len(paths) - len(pptx)
        if skipped:
            messagebox.showwarning(
                "Unsupported files",
                f"Only .pptx files are supported. "
                f"{skipped} file{'s' if skipped != 1 else ''} skipped.",
                parent=self,
            )
        if pptx:
            self._add_files(pptx)
        self._file_panel.set_drop_highlight(False)

    def _on_drop_enter(self) -> None:
        self._file_panel.set_drop_highlight(True)

    def _on_drop_leave(self) -> None:
        self._file_panel.set_drop_highlight(False)

    # ------------------------------------------------------------------
    # File management
    # ------------------------------------------------------------------

    def _browse_pptx(self) -> None:
        paths = filedialog.askopenfilenames(
            title="Select PowerPoint file(s)",
            filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")],
        )
        if paths:
            self._add_files(list(paths))

    def _add_files(self, new_paths: list[str]) -> None:
        """Validate and append new .pptx files. Warns on duplicates."""
        dupes = [p for p in new_paths if p in self._pptx_paths]
        if dupes:
            names = ", ".join(Path(p).name for p in dupes)
            messagebox.showwarning(
                "Duplicate files",
                f"Already added: {names}",
                parent=self,
            )
        unique = [p for p in new_paths if p not in self._pptx_paths]
        if not unique:
            return

        # Read slide counts for new files
        for path in unique:
            try:
                from pptx import Presentation
                self._slide_counts[path] = len(Presentation(path).slides)
            except Exception:
                self._slide_counts[path] = 0

        self._pptx_paths.extend(unique)
        logger.debug("Added %d file(s), total now %d",
                     len(unique), len(self._pptx_paths))

        # Set smart output default if not yet set
        if self._output_dir is None:
            if len(self._pptx_paths) == 1:
                p = Path(self._pptx_paths[0])
                default_out = str(p.parent / f"{p.stem}_pngs")
            else:
                default_out = str(Path(self._pptx_paths[0]).parent)
            self._output_dir = default_out
            self._out_var.set(default_out)

        self._refresh_file_panel()
        self._update_slides_section()
        self._update_run_button_state()

    def _remove_file(self, path: str) -> None:
        """Remove a single file from the list."""
        if path in self._pptx_paths:
            self._pptx_paths.remove(path)
            self._slide_counts.pop(path, None)
        if not self._pptx_paths:
            self._clear_pptx()
            return
        self._refresh_file_panel()
        self._update_slides_section()
        self._update_run_button_state()

    def _clear_pptx(self) -> None:
        self._pptx_paths = []
        self._slide_counts.clear()
        self._file_panel.show_empty()
        self._slides_frame.grid_remove()
        self._update_run_button_state()

    def _refresh_file_panel(self) -> None:
        """Rebuild the file panel display from current state."""
        file_infos = [(p, self._slide_counts.get(p, 0))
                      for p in self._pptx_paths]
        self._file_panel.set_files(file_infos)

    def _update_slides_section(self) -> None:
        """Show slide selection only for exactly 1 file with >1 slide."""
        if (len(self._pptx_paths) == 1
                and self._slide_counts.get(self._pptx_paths[0], 0) > 1):
            self._slides_frame.grid()
            self._all_slides_var.set(True)
            self._slide_range_entry.grid_remove()
        else:
            self._slides_frame.grid_remove()

    def _on_slide_selection_toggle(self) -> None:
        if self._all_slides_var.get():
            self._slide_range_entry.grid_remove()
        else:
            self._slide_range_entry.grid()
            self._slide_range_entry.focus_set()

    def _browse_output(self) -> None:
        path = filedialog.askdirectory(title="Select output folder")
        if path:
            self._output_dir = path
            self._out_var.set(path)
            logger.debug("Selected output dir: %s", path)
            _save_settings({"ppi": self._ppi, "output_dir": path})
            self._update_run_button_state()

    def _on_ppi_change(self, label: str) -> None:
        if label == "Custom":
            self._custom_ppi_frame.grid()
            self._custom_ppi_entry.delete(0, tk.END)
            self._custom_ppi_entry.insert(0, str(self._ppi))
            self._custom_ppi_entry.focus_set()
            self._custom_ppi_entry.select_range(0, tk.END)
            return
        self._custom_ppi_frame.grid_remove()
        self._ppi = _PPI_PRESETS[label]
        self._save_ppi()

    def _apply_custom_ppi(self) -> None:
        raw = self._custom_ppi_entry.get().strip()
        try:
            value = int(raw)
        except ValueError:
            return
        value = max(_PPI_MIN, min(_PPI_MAX, value))
        self._custom_ppi_entry.delete(0, tk.END)
        self._custom_ppi_entry.insert(0, str(value))
        self._ppi = value
        self._save_ppi()

    def _save_ppi(self) -> None:
        logger.debug("PPI set to %d", self._ppi)
        _save_settings({"ppi": self._ppi, "output_dir": self._output_dir})

    def _on_run(self) -> None:
        if not self._pptx_paths or not self._output_dir:
            return
        # Parse slide selection (single-file only).
        self._slide_indices = None
        sc = self._slide_counts.get(self._pptx_paths[0], 0)
        if (len(self._pptx_paths) == 1
                and sc > 1
                and not self._all_slides_var.get()):
            spec = self._slide_range_entry.get().strip()
            if not spec:
                messagebox.showwarning(
                    "No slides specified",
                    "Enter a slide range (e.g. 1-5, 8) or check \"All slides\".",
                    parent=self,
                )
                return
            try:
                self._slide_indices = parse_slide_range(spec, sc)
            except ValueError as exc:
                messagebox.showwarning("Invalid slide range", str(exc), parent=self)
                return
        # Warn if the output folder already contains slide PNGs.
        out = Path(self._output_dir)
        glob_pattern = ("**/slide_*.png" if len(self._pptx_paths) > 1
                        else "slide_*.png")
        if out.is_dir() and list(out.glob(glob_pattern)):
            ok = messagebox.askyesno(
                "Overwrite existing files?",
                f"The output folder already contains exported slides.\n\n"
                f"{self._output_dir}\n\n"
                "Existing slide PNGs will be overwritten. Continue?",
                parent=self,
            )
            if not ok:
                return
        self._error_banner.grid_remove()
        self._cancel_event = threading.Event()
        self._set_ui_busy(True)
        thread = threading.Thread(
            target=self._run_export, daemon=True
        )
        thread.start()

    def _on_cancel(self) -> None:
        if self._cancel_event:
            self._cancel_event.set()
        self._cancel_btn.configure(state="disabled", text="Cancelling…")
        self._status_var.set("Cancelling — finishing current slide…")

    def _open_output_folder(self) -> None:
        if not self._output_dir:
            return
        if sys.platform == "darwin":
            subprocess.run(["open", self._output_dir], check=False)
        elif sys.platform == "win32":
            os.startfile(self._output_dir)  # type: ignore[attr-defined]

    # ------------------------------------------------------------------
    # Export thread
    # ------------------------------------------------------------------

    def _run_export(self) -> None:
        try:
            paths = list(self._pptx_paths)
            if len(paths) == 1:
                self._exporter.export(
                    paths[0],
                    self._output_dir,
                    progress_callback=self._on_progress,
                    cancel_event=self._cancel_event,
                    ppi=self._ppi,
                    slide_indices=self._slide_indices,
                )
            else:
                self._run_batch_export(paths)
            self.after(0, self._on_export_done)
        except InterruptedError:
            self.after(0, self._on_export_cancelled)
        except Exception as exc:
            logger.exception("Export failed")
            msg = str(exc)
            if self._output_dir:
                exported = list(Path(self._output_dir).glob(
                    "**/slide_*.png" if len(self._pptx_paths) > 1
                    else "slide_*.png"
                ))
                if exported:
                    msg += (f"\n\n{len(exported)} slide(s) were exported "
                            "before the error.")
            self.after(0, self._on_export_error, msg)

    def _run_batch_export(self, paths: list[str]) -> None:
        """Export multiple .pptx files, each into its own subfolder."""
        from pptx import Presentation as PptxPresentation

        slide_counts = []
        for p in paths:
            try:
                slide_counts.append(len(PptxPresentation(p).slides))
            except Exception:
                slide_counts.append(1)
        grand_total = sum(slide_counts)

        for file_idx, pptx_path in enumerate(paths):
            if self._cancel_event and self._cancel_event.is_set():
                raise InterruptedError("Export cancelled by user.")

            stem = Path(pptx_path).stem
            file_out = str(Path(self._output_dir) / f"{stem}_pngs")
            offset = sum(slide_counts[:file_idx])

            def _make_cb(off, gt, fidx, fname, fcount):
                def cb(current, total):
                    abs_current = off + current
                    frac = abs_current / gt if gt > 0 else 0
                    if current < total:
                        msg = (f"[{fidx}/{fcount}] {fname}: "
                               f"slide {current + 1}/{total}")
                    else:
                        msg = f"[{fidx}/{fcount}] {fname}: done"
                    self.after(0, self._update_progress, frac, msg)
                return cb

            self._exporter.export(
                pptx_path,
                file_out,
                progress_callback=_make_cb(
                    offset, grand_total, file_idx + 1, stem, len(paths)),
                cancel_event=self._cancel_event,
                ppi=self._ppi,
            )

    def _on_progress(self, current: int, total: int) -> None:
        if total == 0:
            return
        fraction = current / total
        msg = (
            f"Processing slide {current + 1} of {total}…"
            if current < total
            else "Finalising…"
        )
        self.after(0, self._update_progress, fraction, msg)

    def _update_progress(self, fraction: float, msg: str) -> None:
        self._progress_bar.set(fraction)
        self._status_var.set(msg)

    def _on_export_done(self) -> None:
        self._progress_bar.set(1.0)
        n = len(self._pptx_paths)
        if n > 1:
            self._status_var.set(f"Done — {n} presentations exported.")
        else:
            self._status_var.set("Done — all slides exported.")
        self._set_ui_busy(False)
        self._open_folder_btn.grid()

    def _on_export_cancelled(self) -> None:
        self._progress_bar.set(0)
        self._status_var.set("Export cancelled.")
        self._set_ui_busy(False)

    def _on_export_error(self, message: str) -> None:
        self._set_ui_busy(False)
        self._error_banner.show(message)
        self._error_banner.grid()

    def _dismiss_error(self) -> None:
        self._error_banner.grid_remove()
        self._status_var.set("")

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _update_run_button_state(self) -> None:
        ready = bool(self._powerpoint_available and self._pptx_paths and self._output_dir)
        self._run_btn.configure(state="normal" if ready else "disabled")

    def _set_ui_busy(self, busy: bool) -> None:
        if busy:
            self._progress_bar.set(0)
            self._status_var.set("Starting export…")
            self._open_folder_btn.grid_remove()
            self._progress_frame.grid()
            self._run_btn.grid_remove()
            self._cancel_btn.configure(state="normal", text="Cancel")
            self._cancel_btn.grid()
        else:
            self._cancel_btn.grid_remove()
            self._run_btn.grid()
            self._update_run_button_state()
        self.update_idletasks()


# ---------------------------------------------------------------------------
# Reusable sub-widgets
# ---------------------------------------------------------------------------

class _FilePanel(ctk.CTkFrame):
    """File selection panel with file list, add/remove, and drop support."""

    def __init__(self, parent, on_browse, on_clear_all, on_remove_file,
                 dnd_enabled: bool = False):
        super().__init__(
            parent,
            fg_color=_COLOR_BG_SURFACE,
            border_color=_COLOR_BORDER,
            border_width=1,
            corner_radius=8,
        )
        self.grid_columnconfigure(0, weight=1)
        self._on_browse = on_browse
        self._on_clear_all = on_clear_all
        self._on_remove_file = on_remove_file
        self._dnd_enabled = dnd_enabled
        self._file_rows: dict[str, ctk.CTkFrame] = {}

        self._build_empty()
        self._build_files()
        self._show_state_empty()

    # ── Empty state ────────────────────────────────────────────────────

    def _build_empty(self) -> None:
        self._empty_frame = ctk.CTkFrame(self, fg_color="transparent")
        self._empty_frame.grid_columnconfigure(0, weight=1)

        ctk.CTkLabel(
            self._empty_frame,
            text="No files selected",
            font=_FONT_BODY_BOLD,
            text_color=_COLOR_TEXT_PRIMARY,
        ).grid(row=0, column=0, pady=(16, 2))

        subline = ("Drop .pptx files here or click Browse"
                   if self._dnd_enabled
                   else "Click Browse to open .pptx files")
        ctk.CTkLabel(
            self._empty_frame,
            text=subline,
            font=_FONT_SMALL,
            text_color=_COLOR_TEXT_SECONDARY,
        ).grid(row=1, column=0, pady=(0, 12))

        ctk.CTkButton(
            self._empty_frame,
            text="Browse\u2026",
            command=self._on_browse,
            width=110,
            font=_FONT_BODY,
            fg_color=_COLOR_ACCENT,
            hover_color=_COLOR_ACCENT_HOVER,
        ).grid(row=2, column=0, pady=(0, 16))

    # ── Files-loaded state ─────────────────────────────────────────────

    def _build_files(self) -> None:
        self._files_frame = ctk.CTkFrame(self, fg_color="transparent")
        self._files_frame.grid_columnconfigure(0, weight=1)

        # Scrollable list area for file rows
        self._list_frame = ctk.CTkFrame(self._files_frame, fg_color="transparent")
        self._list_frame.grid(row=0, column=0, sticky="ew", padx=12, pady=(10, 0))
        self._list_frame.grid_columnconfigure(0, weight=1)

        # Separator
        ctk.CTkFrame(
            self._files_frame, height=1, fg_color=_COLOR_BORDER, corner_radius=0,
        ).grid(row=1, column=0, sticky="ew", padx=12, pady=(8, 0))

        # Bottom bar: + Add files  |  Clear all
        btn_bar = ctk.CTkFrame(self._files_frame, fg_color="transparent")
        btn_bar.grid(row=2, column=0, sticky="ew", padx=12, pady=(6, 10))
        btn_bar.grid_columnconfigure(0, weight=1)

        ctk.CTkButton(
            btn_bar,
            text="+ Add files",
            command=self._on_browse,
            width=100,
            height=28,
            font=_FONT_SMALL,
            fg_color="transparent",
            text_color=_COLOR_ACCENT,
            hover_color=_COLOR_BG_BASE,
            border_width=1,
            border_color=_COLOR_ACCENT,
            corner_radius=4,
        ).grid(row=0, column=0, sticky="w")

        ctk.CTkButton(
            btn_bar,
            text="Clear all",
            command=self._on_clear_all,
            width=80,
            height=28,
            font=_FONT_SMALL,
            fg_color="transparent",
            text_color=_COLOR_TEXT_SECONDARY,
            hover_color=_COLOR_BG_BASE,
            border_width=1,
            border_color=_COLOR_BORDER,
            corner_radius=4,
        ).grid(row=0, column=1, sticky="e")

    # ── State switching ────────────────────────────────────────────────

    def _show_state_empty(self) -> None:
        self._files_frame.grid_remove()
        self._empty_frame.grid(row=0, column=0, sticky="ew")

    def _show_state_files(self) -> None:
        self._empty_frame.grid_remove()
        self._files_frame.grid(row=0, column=0, sticky="ew")

    # ── Public API ─────────────────────────────────────────────────────

    def set_files(self, file_infos: list[tuple[str, int]]) -> None:
        """Update the displayed file list.

        *file_infos*: list of ``(path, slide_count)`` tuples.
        """
        if not file_infos:
            self.show_empty()
            return

        # Clear existing rows
        for widget in self._file_rows.values():
            widget.destroy()
        self._file_rows.clear()

        for i, (path, slide_count) in enumerate(file_infos):
            self._add_row(i, path, slide_count)

        self._show_state_files()

    def show_empty(self) -> None:
        """Reset to the empty / no-files state."""
        for widget in self._file_rows.values():
            widget.destroy()
        self._file_rows.clear()
        self._show_state_empty()

    def set_drop_highlight(self, active: bool) -> None:
        """Toggle accent border for drag-over visual feedback."""
        self.configure(
            border_color=_COLOR_ACCENT if active else _COLOR_BORDER,
        )

    # ── Row construction ───────────────────────────────────────────────

    def _add_row(self, index: int, path: str, slide_count: int) -> None:
        row = ctk.CTkFrame(self._list_frame, fg_color="transparent")
        row.grid(row=index, column=0, sticky="ew", pady=(0, 2))
        row.grid_columnconfigure(0, weight=1)

        ctk.CTkLabel(
            row,
            text=Path(path).name,
            font=_FONT_BODY,
            text_color=_COLOR_TEXT_PRIMARY,
            anchor="w",
        ).grid(row=0, column=0, sticky="w")

        if slide_count > 0:
            slides_text = f"{slide_count} slide{'s' if slide_count != 1 else ''}"
            ctk.CTkLabel(
                row,
                text=slides_text,
                font=_FONT_SMALL,
                text_color=_COLOR_TEXT_SECONDARY,
                anchor="e",
            ).grid(row=0, column=1, padx=(8, 0))

        p = path  # capture for lambda closure
        ctk.CTkButton(
            row,
            text="\u2715",
            command=lambda: self._on_remove_file(p),
            width=24,
            height=24,
            font=_FONT_SMALL,
            fg_color="transparent",
            text_color=_COLOR_TEXT_SECONDARY,
            hover_color=_COLOR_BORDER,
            corner_radius=4,
        ).grid(row=0, column=2, padx=(4, 0))

        self._file_rows[path] = row


class _ErrorBanner(ctk.CTkFrame):
    """Inline error banner — shown below the run button on failure."""

    def __init__(self, parent, on_dismiss):
        super().__init__(
            parent,
            fg_color=_COLOR_ERROR_BG,
            corner_radius=8,
        )
        self.grid_columnconfigure(0, weight=1)
        self._on_dismiss = on_dismiss

        self._msg_label = ctk.CTkLabel(
            self,
            text="",
            font=_FONT_SMALL,
            text_color=_COLOR_ERROR_TEXT,
            anchor="w",
            wraplength=400,
            justify="left",
        )
        self._msg_label.grid(row=0, column=0, sticky="ew", padx=12, pady=10)

        ctk.CTkButton(
            self,
            text="Dismiss",
            command=self._on_dismiss,
            width=70,
            height=26,
            font=_FONT_SMALL,
            fg_color="transparent",
            text_color=_COLOR_ERROR_TEXT,
            hover_color=_COLOR_ERROR_BG,
            border_width=1,
            border_color=_COLOR_ERROR_TEXT,
            corner_radius=4,
        ).grid(row=0, column=1, padx=(0, 12), pady=10)

    def show(self, message: str) -> None:
        self._msg_label.configure(text=message)
