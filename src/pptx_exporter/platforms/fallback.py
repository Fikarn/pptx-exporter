"""Fallback backend: python-pptx + Pillow compositing.

Used when Microsoft PowerPoint is not installed. Renders each slide's shapes
onto a transparent RGBA canvas at the slide's native resolution, ignoring the
slide background entirely.

Per-slide workflow:
1. Add an invisible full-slide bounding rectangle (no fill, no border).
2. Composite all shapes onto a transparent canvas.
3. Remove the bounding rectangle.
4. Save as PNG.
"""

import logging
from pathlib import Path
from typing import Callable, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Emu

from ..utils import slide_output_name

logger = logging.getLogger(__name__)


def _emu_to_px(emu: int, dpi: int = 96) -> int:
    """Convert EMU (English Metric Units) to pixels at *dpi*."""
    # 1 inch = 914400 EMU; 1 inch = dpi pixels
    return round(emu * dpi / 914400)


def export_slides(
    pptx_path: Path,
    output_dir: Path,
    dpi: int = 96,
    progress_callback: Optional[Callable[[int, int], None]] = None,
) -> None:
    """Export every slide of *pptx_path* as a transparent PNG into *output_dir*.

    Args:
        pptx_path: Resolved path to the .pptx file.
        output_dir: Resolved path to the output directory.
        dpi: Resolution for rasterisation (default 96 dpi = 100 % screen).
        progress_callback: Optional callable(slide_index, total_slides).
    """
    try:
        from PIL import Image
    except ImportError as exc:
        raise RuntimeError(
            "Pillow is required for the fallback backend. "
            "Install it with: pip install Pillow"
        ) from exc

    try:
        from pptx2svg import convert  # noqa: F401 — optional, not used yet
    except ImportError:
        pass  # pptx2svg is optional

    prs = Presentation(str(pptx_path))
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    slide_w_px = _emu_to_px(slide_width_emu, dpi)
    slide_h_px = _emu_to_px(slide_height_emu, dpi)

    total = len(prs.slides)

    logger.info(
        "Fallback export: %d slides, %dx%d px @ %d dpi",
        total,
        slide_w_px,
        slide_h_px,
        dpi,
    )

    for idx, slide in enumerate(prs.slides):
        if progress_callback:
            progress_callback(idx, total)

        logger.debug("Processing slide %d/%d", idx + 1, total)

        # Step 1 — add invisible bounding rectangle
        bounding_shape = _add_bounding_rect(slide, slide_width_emu, slide_height_emu)

        # Step 2+3 — composite shapes onto transparent canvas
        canvas = Image.new("RGBA", (slide_w_px, slide_h_px), (0, 0, 0, 0))
        canvas = _composite_shapes(canvas, slide, slide_w_px, slide_h_px, dpi)

        # Step 4 — remove bounding rectangle
        _remove_shape(slide, bounding_shape)

        out_name = slide_output_name(idx, total)
        out_path = output_dir / out_name
        canvas.save(str(out_path), "PNG")
        logger.info("Saved %s", out_path)

    if progress_callback:
        progress_callback(total, total)


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _add_bounding_rect(slide, width_emu: int, height_emu: int):
    """Add an invisible full-slide rectangle to *slide* and return it."""
    from pptx.util import Emu as _Emu

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE == 1; avoids importing MSO_AUTO_SHAPE_TYPE
        _Emu(0),
        _Emu(0),
        _Emu(width_emu),
        _Emu(height_emu),
    )
    # No fill
    shape.fill.background()
    # No line
    shape.line.fill.background()
    return shape


def _remove_shape(slide, shape) -> None:
    """Remove *shape* from *slide*'s shape tree."""
    sp = shape._element
    sp.getparent().remove(sp)


def _composite_shapes(canvas, slide, slide_w_px: int, slide_h_px: int, dpi: int):
    """Render each shape in *slide* onto *canvas*.

    Only handles picture/image shapes with available image blobs. Non-image
    shapes (text boxes, auto-shapes) are skipped with a debug log — rendering
    arbitrary vector shapes correctly requires a full SVG renderer which is
    beyond the scope of the fallback path.
    """
    from PIL import Image

    for shape in slide.shapes:
        if not shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE == 13
            logger.debug("Skipping non-picture shape: %s (type %s)", shape.name, shape.shape_type)
            continue
        try:
            import io

            img_data = shape.image.blob
            img = Image.open(io.BytesIO(img_data)).convert("RGBA")

            left_px = _emu_to_px(shape.left, dpi)
            top_px = _emu_to_px(shape.top, dpi)
            width_px = _emu_to_px(shape.width, dpi)
            height_px = _emu_to_px(shape.height, dpi)

            img = img.resize((width_px, height_px), Image.LANCZOS)
            canvas.paste(img, (left_px, top_px), img)
            logger.debug("Composited picture shape: %s", shape.name)
        except Exception as exc:
            logger.warning("Could not render shape '%s': %s", shape.name, exc)

    return canvas
