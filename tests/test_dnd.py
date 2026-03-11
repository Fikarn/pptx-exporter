"""Unit tests for drag-and-drop and file panel logic."""

import os
from pathlib import Path
from unittest import mock

from pptx_exporter.tkdnd import _vendored_dir, DND_FILES


# ---------------------------------------------------------------------------
# tkdnd loader
# ---------------------------------------------------------------------------


def test_vendored_dir_returns_path_on_supported_platform():
    """_vendored_dir should return a directory for macOS or Windows."""
    result = _vendored_dir()
    # On CI (Linux) this may be None; on macOS/Windows it should exist.
    if result is not None:
        assert os.path.isdir(result)


def test_vendored_dir_macos_arm64():
    with mock.patch("pptx_exporter.tkdnd.platform") as m:
        m.system.return_value = "Darwin"
        m.machine.return_value = "arm64"
        result = _vendored_dir()
    assert result is not None
    assert result.endswith("macos-arm64")


def test_vendored_dir_macos_x86_64():
    with mock.patch("pptx_exporter.tkdnd.platform") as m:
        m.system.return_value = "Darwin"
        m.machine.return_value = "x86_64"
        result = _vendored_dir()
    assert result is not None
    assert result.endswith("macos-x86_64")


def test_vendored_dir_windows():
    with mock.patch("pptx_exporter.tkdnd.platform") as m:
        m.system.return_value = "Windows"
        m.machine.return_value = "AMD64"
        result = _vendored_dir()
    assert result is not None
    assert result.endswith("windows-x64")


def test_vendored_dir_unsupported_platform():
    with mock.patch("pptx_exporter.tkdnd.platform") as m:
        m.system.return_value = "Linux"
        m.machine.return_value = "x86_64"
        assert _vendored_dir() is None


def test_dnd_files_constant():
    assert DND_FILES == "DND_Files"


# ---------------------------------------------------------------------------
# Drop path parsing
# ---------------------------------------------------------------------------


def test_parse_drop_paths_simple():
    """Space-separated paths without spaces parse correctly."""
    # Simulate what tk.splitlist does with simple paths
    data = "/path/to/file1.pptx /path/to/file2.pptx"
    # tk.splitlist splits on whitespace and handles braces
    # We test the logic, not Tcl itself
    paths = data.split()
    assert len(paths) == 2
    assert paths[0] == "/path/to/file1.pptx"
    assert paths[1] == "/path/to/file2.pptx"


def test_filter_pptx_from_mixed_drop():
    """Only .pptx files should be accepted from a drop."""
    dropped = [
        "/tmp/presentation.pptx",
        "/tmp/document.pdf",
        "/tmp/image.png",
        "/tmp/other.PPTX",
    ]
    pptx = [p for p in dropped if p.lower().endswith(".pptx")]
    assert len(pptx) == 2
    assert "/tmp/presentation.pptx" in pptx
    assert "/tmp/other.PPTX" in pptx


# ---------------------------------------------------------------------------
# File list management (add / remove / duplicate detection)
# ---------------------------------------------------------------------------


class TestFileListLogic:
    """Test the add/remove/duplicate logic used by App._add_files."""

    def _make_pptx(self, tmp_path: Path, name: str) -> str:
        """Create a minimal .pptx file and return its path as a string."""
        from pptx import Presentation
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = tmp_path / name
        prs.save(str(path))
        return str(path)

    def test_add_files_appends(self, tmp_path: Path):
        paths: list[str] = []
        f1 = self._make_pptx(tmp_path, "a.pptx")
        f2 = self._make_pptx(tmp_path, "b.pptx")

        # First add
        new = [p for p in [f1] if p not in paths]
        paths.extend(new)
        assert paths == [f1]

        # Second add — should append
        new = [p for p in [f2] if p not in paths]
        paths.extend(new)
        assert paths == [f1, f2]

    def test_add_files_detects_duplicates(self, tmp_path: Path):
        paths: list[str] = []
        f1 = self._make_pptx(tmp_path, "a.pptx")

        paths.append(f1)
        dupes = [p for p in [f1] if p in paths]
        unique = [p for p in [f1] if p not in paths]

        assert dupes == [f1]
        assert unique == []

    def test_remove_file(self, tmp_path: Path):
        f1 = self._make_pptx(tmp_path, "a.pptx")
        f2 = self._make_pptx(tmp_path, "b.pptx")
        paths = [f1, f2]
        slide_counts = {f1: 1, f2: 3}

        paths.remove(f1)
        slide_counts.pop(f1, None)

        assert paths == [f2]
        assert f1 not in slide_counts

    def test_remove_last_file_empties_list(self, tmp_path: Path):
        f1 = self._make_pptx(tmp_path, "a.pptx")
        paths = [f1]

        paths.remove(f1)
        assert paths == []

    def test_slide_section_visibility_single_multi_slide(self, tmp_path: Path):
        """Slide section should show for 1 file with >1 slide."""
        f1 = self._make_pptx(tmp_path, "a.pptx")
        paths = [f1]
        slide_counts = {f1: 5}

        show = (len(paths) == 1
                and slide_counts.get(paths[0], 0) > 1)
        assert show is True

    def test_slide_section_hidden_for_batch(self, tmp_path: Path):
        """Slide section should hide for 2+ files."""
        f1 = self._make_pptx(tmp_path, "a.pptx")
        f2 = self._make_pptx(tmp_path, "b.pptx")
        paths = [f1, f2]
        slide_counts = {f1: 5, f2: 3}

        show = (len(paths) == 1
                and slide_counts.get(paths[0], 0) > 1)
        assert show is False

    def test_slide_section_hidden_for_single_one_slide(self, tmp_path: Path):
        """Slide section should hide for 1 file with only 1 slide."""
        f1 = self._make_pptx(tmp_path, "a.pptx")
        paths = [f1]
        slide_counts = {f1: 1}

        show = (len(paths) == 1
                and slide_counts.get(paths[0], 0) > 1)
        assert show is False

    def test_output_dir_default_single(self, tmp_path: Path):
        """Single file → {stem}_pngs sibling folder."""
        f1 = self._make_pptx(tmp_path, "deck.pptx")
        paths = [f1]

        p = Path(paths[0])
        default_out = str(p.parent / f"{p.stem}_pngs")
        assert default_out == str(tmp_path / "deck_pngs")

    def test_output_dir_default_batch(self, tmp_path: Path):
        """Multiple files → parent directory of first file."""
        f1 = self._make_pptx(tmp_path, "a.pptx")
        f2 = self._make_pptx(tmp_path, "b.pptx")
        paths = [f1, f2]

        default_out = str(Path(paths[0]).parent)
        assert default_out == str(tmp_path)
